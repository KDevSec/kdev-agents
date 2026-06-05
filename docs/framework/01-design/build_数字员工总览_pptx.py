# -*- coding: utf-8 -*-
"""手搓 python-pptx，按 1280x720 复刻 KDev 数字员工概念模型汇报原型（单页总览）。
竖向铺开填满整页（内容 y=74..690），避免上紧下空。
输出：KDev数字员工概念模型_汇报原型_v0.1.pptx"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "KDev数字员工概念模型_汇报原型_v0.1.pptx")
FONT = "Microsoft YaHei"

NAVY="1F4E79"; NAVYD="1F3864"; RED="C00000"; BBLUE="2E75B6"
LT="F2F6FB"; BD="D9D9D9"; CELL="BFBFBF"; INK="333333"; MUT="666666"
WARN="C55A11"; OKG="2E7D32"; WHITE="FFFFFF"

EMU_PX = 9525
def PX(v): return Emu(int(round(v*EMU_PX)))

prs = Presentation()
prs.slide_width = PX(1280); prs.slide_height = PX(720)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(x,y,w,h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, PX(x),PX(y),PX(w),PX(h))
    s.shadow.inherit=False
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(fill)
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=RGBColor.from_string(line); s.line.width=Pt(lw)
    return s

def text(x,y,w,h, segs, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(PX(x),PX(y),PX(w),PX(h))
    tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Emu(0); tf.margin_right=Emu(0); tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    if isinstance(segs,str): segs=[(segs,color,bold,size)]
    p = tf.paragraphs[0]; p.alignment=align
    if line_spacing: p.line_spacing=line_spacing
    for seg in segs:
        t=seg[0]; c=seg[1] if len(seg)>1 else color; b=seg[2] if len(seg)>2 else bold
        sz=seg[3] if len(seg)>3 else size
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b
        r.font.name=FONT; r.font.color.rgb=RGBColor.from_string(c)
        rPr=r._r.get_or_add_rPr()
        latin=rPr.find(qn('a:latin')); ea=rPr.find(qn('a:ea'))
        if ea is None:
            ea=rPr.makeelement(qn('a:ea'),{})
            (latin.addnext(ea) if latin is not None else rPr.append(ea))
        ea.set('typeface', FONT)
    return tb

def panel(x,y,w,h, title, hc=NAVY):
    rect(x,y,w,h, WHITE, BD, 1)
    rect(x,y,w,28, hc, None)
    text(x+11,y,w-22,28, title, 13.5, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    return x, y+28

# ===== 顶部标题 =====
rect(32,18,34,34, NAVYD, None)
text(32,18,34,34, "1", 21, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(79,18,46,34, "总览", 17, RED, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
text(132,16,600,38, "KDev 数字员工 · 架构总览", 23, NAVY, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
rect(32,60,1216,3, NAVY, None)

# ============ 左栏 x=32 w=675 ============
LX=32; LW=675; cx=LX+LW/2
# --- Panel 1：组织架构 (74..374) ---
panel(LX,74,LW,300, " 数字员工 = 多能力 ＋ 编排 → 6 人 AI 团队（39 agent/能力）")
rect(cx-34,118,68,24, WHITE, BD, 1); text(cx-34,118,68,24,"用户",12,INK,False,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
rect(cx-1,142,2,14, BD, None)
rect(cx-180,160,360,30, NAVY, None)
text(cx-180,160,360,30,[("CEO 主控员  ",WHITE,True,13),("对话/派单/决策/升级裁决",WHITE,False,10.5)],13,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
rect(cx-1,190,2,14, BD, None)
rect(44,208,110,104, RED, None)
text(44,208,110,104,[("CQO 元监督\n",WHITE,True,12.5),("流程合规·有阻断权",WHITE,False,10)],12.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
rect(164,208,531,104, WHITE, BBLUE, 1.5)
text(171,213,517,18,"业务员工（每个 = 一个\"数字员工\"）＋ 评审专家收敛把关",11.5,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
emps=[("需求架构师","5 能力"),("开发工程师","6 能力"),("测试工程师","4 能力")]; ew=135
for i,(en,ec) in enumerate(emps):
    ex=171+i*(ew+6)
    rect(ex,236,ew,52, LT, BBLUE, 1)
    text(ex,242,ew,18,en,11.5,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
    text(ex,264,ew,16,ec,10.5,MUT,False,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
rect(594,236,92,52, NAVY, NAVY, 1.5)
text(594,242,92,18,"评审专家",11.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
text(594,264,92,16,"18 评审·收敛",9.5,"DBE6F2",False,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
text(171,294,517,14,"↕ 3 交付员工编排间直接协作（不经 CEO）· 产出 → 评审专家收敛评审",9.5,MUT,False,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(44,324,651,40,[("39 agent",NAVY,True,11.5),(" = 1 CEO ＋ 1 CQO ＋ 4 编排 ＋ 15 业务 ＋ 18 评审\n",INK,False,11.5),("评审治理",NAVY,True,11.5),("：R→D→T→F（19 节点）· 双重通过门（总分≥阈值 且 阻断=0，3 次不过升 CEO）",INK,False,11.5)],11.5,INK,False,PP_ALIGN.LEFT,MSO_ANCHOR.TOP,line_spacing=1.4)

# --- Panel 2：金字塔 (388..688) ---
panel(LX,388,LW,300, "从底层技能到数字员工（金字塔）")
pyr=[("数字员工","编排 + 多能力 → 端到端交付",0.50,NAVYD),
     ("能力（= 专用 agent）","某类活的专家",0.68,NAVY),
     ("底层技能：Skill ＋ 规范 ＋ 模板 ＋ 知识库","",0.86,"2E5C8A"),
     ("通用 Agent（底座 · 推理 + 工具/MCP + 模型选型）","",1.0,BBLUE)]
inner=653; py=426
for label,sub,frac,col in pyr:
    pw=inner*frac; pxx=cx-pw/2
    rect(pxx,py,pw,48,col,None)
    if sub:
        text(pxx,py,pw,48,[(label+"  ",WHITE,True,12.5),(sub,"DBE6F2",False,10)],12.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    else:
        text(pxx,py,pw,48,label,12.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    py+=56
text(LX,py+4,LW,18,"自下而上逐层组装：底座 → 底层技能 → 能力 → 数字员工",10.5,MUT,False,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)

# ============ 右栏 x=720 w=528 ============
RX=720; RW=528
# --- Panel A：能力组成 (74..254) ---
panel(RX,74,RW,180, "① 能力组成 ＋ 平台底座")
caps=[("通用 Agent","推理+工具/MCP"),("Skill","专业 know-how"),("规范","质量约束"),("模板","产出结构")]; cw=122
for i,(t1,t2) in enumerate(caps):
    cxx=RX+11+i*(cw+6)
    rect(cxx,116,cw,52, WHITE, BBLUE, 1.5)
    text(cxx,122,cw,18,t1,11.5,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
    text(cxx,142,cw,16,t2,9.5,MUT,False,PP_ALIGN.CENTER,MSO_ANCHOR.TOP)
text(RX+11,180,RW-22,18,[("专用agent",NAVY,True,11.5),("：由上面四项组成专用Agent，比如测试用例 agent / 安全编码 agent",INK,False,11.5)],11.5,INK,False,PP_ALIGN.LEFT,MSO_ANCHOR.TOP)
text(RX+11,212,RW-22,18,[("平台底座",NAVY,True,11.5),("：kdev-staff / kdev-core / kdev-memory",INK,False,11.5)],11.5,INK,False,PP_ALIGN.LEFT,MSO_ANCHOR.TOP)

# --- Panel B：自主度表 (270..486) ---
panel(RX,270,RW,216, "② 自主度分级")
rows=[("级","名","人/机分工","按级长出"),
      ("L1","助手","人驱动、执行单步","能力起步"),
      ("L2","协同","员工做完整任务、人编排确认","＋编排＋记忆"),
      ("L3","自主","端到端自跑、人管异常","＋协作＋第三方评审"),
      ("L4","自治","自监督/自演进","＋元监督＋自演进")]
tx=RX+11; ty=306; tw=RW-22
colw=[int(tw*0.10),int(tw*0.18),int(tw*0.40)]; colw.append(tw-sum(colw)); rh=31
for ri,row in enumerate(rows):
    cx2=tx
    for ci,val in enumerate(row):
        fill = NAVY if ri==0 else (LT if ri%2==0 else WHITE)
        rect(cx2,ty+ri*rh,colw[ci],rh, fill, CELL, 0.75)
        if ri==0:
            text(cx2+5,ty+ri*rh,colw[ci]-7,rh,val,11,WHITE,True,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
        else:
            c=INK; b=False
            if ci==0: c=NAVY; b=True
            if ci==3 and ri==3: c=WARN; b=True
            if ci==3 and ri==4: c=RED; b=True
            text(cx2+5,ty+ri*rh,colw[ci]-7,rh,val,10.5,c,b,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
        cx2+=colw[ci]
text(RX+11, ty+5*rh+6, RW-22,16,"知识库/MCP 不在此轴（属组成与持续升级）。",10.5,RED,True,PP_ALIGN.LEFT,MSO_ANCHOR.TOP)

# --- Panel C：飞轮 (502..688) ---
panel(RX,502,RW,186, "③ 持续升级 · 复利飞轮", hc=RED)
fly=[("1","数字员工干活（需求→验收）"),
     ("2","过程沉淀进 专有知识库（持续注入）"),
     ("3","知识 反哺能力（下次更准更快）"),
     ("4","越用越强（复利上升） → 回到 1")]
fy=540
for n,t in fly:
    rect(RX+12,fy,22,22, NAVY, None, shape=MSO_SHAPE.OVAL)
    text(RX+12,fy,22,22,n,11.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    col = RED if n=="4" else INK; bold=(n=="4")
    text(RX+42,fy,RW-54,22,t,11.5,col,bold,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    fy+=30
text(RX+11,fy+4,RW-22,18,"知识库 md·wiki 或 RAG 皆可、不分高低；是持续升级引擎，不挂级别。",10.5,RED,True,PP_ALIGN.LEFT,MSO_ANCHOR.TOP)

# ===== 页脚 =====
rect(32,701,1216,0.8, "E6E9EE", None)
text(32,703,600,14,"KDev 数字员工架构 · 概念模型",11.5,"9AA3AD",False,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
text(648,703,600,14,"01",11.5,"9AA3AD",False,PP_ALIGN.RIGHT,MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("OK", OUT)
